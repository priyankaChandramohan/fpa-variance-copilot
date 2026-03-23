"""
app.py — FP&A Variance Copilot  ·  Streamlit application entry point.

Three UI states
---------------
State 1  No file uploaded   → Landing page explaining the workflow.
State 2  File uploaded       → KPI cards + four chart tabs.
State 3  Commentary ready    → Everything from State 2 + AI narrative + PPTX download.

Session-state keys used
-----------------------
_file_key    : str   lightweight cache-bust key (filename + size + thresholds)
enriched_df  : DataFrame   output of compute_variances()
cat_df       : DataFrame   output of summarize_by_category()
waterfall_df : DataFrame   output of build_waterfall_data()
commentary   : str | None  Claude response text, None until generated
"""

from __future__ import annotations

import pandas as pd
import streamlit as st
from io import BytesIO

import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from variance_engine import (
    validate_csv,
    compute_variances,
    summarize_by_category,
    build_waterfall_data,
    get_material_items,
)
from charts import (
    plot_waterfall,
    plot_budget_vs_actual,
    plot_severity_donut,
    plot_projection,
)
from commentary import (
    generate_commentary,
    answer_question,
    TONES,
    DEFAULT_TONE,
)


# ── Page config  (must be first Streamlit call) ───────────────────────────────

st.set_page_config(
    page_title="FP&A Variance Copilot",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
)


# ── Custom CSS ────────────────────────────────────────────────────────────────

st.markdown("""
<style>
/* ─── Sidebar dark theme ──────────────────────────────────────────── */
section[data-testid="stSidebar"] {
    background-color: #0F172A;
}
section[data-testid="stSidebar"] .stMarkdown p,
section[data-testid="stSidebar"] .stMarkdown span,
section[data-testid="stSidebar"] label,
section[data-testid="stSidebar"] .stTextInput label,
section[data-testid="stSidebar"] .stSlider label,
section[data-testid="stSidebar"] .stFileUploader label,
section[data-testid="stSidebar"] small {
    color: #94A3B8 !important;
    font-family: 'Calibri', 'Segoe UI', sans-serif;
    font-size: 0.78rem;
    letter-spacing: 0.04em;
    text-transform: uppercase;
}
section[data-testid="stSidebar"] h2 {
    color: #F1F5F9 !important;
    font-family: 'Calibri', 'Segoe UI', sans-serif;
    font-size: 1.1rem;
    font-weight: 700;
    letter-spacing: -0.01em;
    text-transform: none;
}
section[data-testid="stSidebar"] hr {
    border-color: #334155;
    margin: 12px 0;
}
/* Slider track colour */
section[data-testid="stSidebar"] .stSlider [data-testid="stTickBar"] {
    background: #334155;
}

/* ─── KPI cards ───────────────────────────────────────────────────── */
.kpi-card {
    background: #FFFFFF;
    border: 1px solid #E2E8F0;
    border-radius: 8px;
    padding: 20px 16px 16px;
    text-align: center;
    min-height: 110px;
    box-shadow: 0 1px 3px rgba(0,0,0,0.04);
}
.kpi-label {
    font-family: 'Calibri', 'Segoe UI', sans-serif;
    font-size: 0.70rem;
    font-weight: 700;
    color: #64748B;
    letter-spacing: 0.07em;
    text-transform: uppercase;
    margin-bottom: 6px;
}
.kpi-value {
    font-family: 'Calibri', 'Segoe UI', sans-serif;
    font-size: 1.65rem;
    font-weight: 700;
    color: #1E293B;
    line-height: 1.15;
    margin-bottom: 4px;
}
.kpi-value.favorable   { color: #16A34A; }
.kpi-value.unfavorable { color: #DC2626; }
.kpi-subtitle {
    font-family: 'Calibri', 'Segoe UI', sans-serif;
    font-size: 0.70rem;
    color: #94A3B8;
    letter-spacing: 0.02em;
}

/* ─── Landing steps ───────────────────────────────────────────────── */
.landing-step {
    background: #F8FAFC;
    border: 1px solid #E2E8F0;
    border-radius: 12px;
    padding: 32px 24px 28px;
    text-align: center;
    height: 100%;
    transition: box-shadow 0.2s;
}
.landing-step:hover {
    box-shadow: 0 4px 16px rgba(0,0,0,0.06);
}
.step-icon { font-size: 2.2rem; margin-bottom: 14px; }
.landing-step h3 {
    font-family: 'Calibri', 'Segoe UI', sans-serif;
    font-size: 1.05rem;
    font-weight: 700;
    color: #1E293B;
    margin-bottom: 10px;
}
.landing-step p {
    font-family: 'Calibri', 'Segoe UI', sans-serif;
    font-size: 0.875rem;
    color: #64748B;
    line-height: 1.6;
    margin: 0;
}
.landing-connector {
    display: flex;
    align-items: center;
    justify-content: center;
    font-size: 1.4rem;
    color: #CBD5E1;
    padding-top: 40px;
}

/* ─── Commentary box ──────────────────────────────────────────────── */
.commentary-box {
    background: #F8FAFC;
    border: 1px solid #E2E8F0;
    border-left: 4px solid #0D9488;
    border-radius: 0 8px 8px 0;
    padding: 24px 28px;
    font-family: 'Calibri', 'Segoe UI', sans-serif;
    font-size: 0.92rem;
    color: #334155;
    line-height: 1.75;
    white-space: pre-wrap;
}

/* ─── Tabs ────────────────────────────────────────────────────────── */
.stTabs [data-baseweb="tab-list"] {
    gap: 4px;
    border-bottom: 2px solid #E2E8F0;
    padding-bottom: 0;
}
.stTabs [data-baseweb="tab"] {
    font-family: 'Calibri', 'Segoe UI', sans-serif;
    font-size: 0.875rem;
    font-weight: 600;
    color: #64748B;
    padding: 10px 18px;
    border-radius: 6px 6px 0 0;
    background: transparent;
}
.stTabs [aria-selected="true"] {
    color: #0D9488 !important;
    background: #F0FDFA !important;
    border-bottom: 2px solid #0D9488;
}

/* ─── Global typography ───────────────────────────────────────────── */
html, body, [class*="css"], .stMarkdown {
    font-family: 'Calibri', 'Segoe UI', Arial, sans-serif;
}
h1, h2, h3 { color: #1E293B; }

/* ─── Primary button colour override ─────────────────────────────── */
.stButton [data-testid="baseButton-primary"] {
    background-color: #0D9488;
    border-color: #0D9488;
}
.stButton [data-testid="baseButton-primary"]:hover {
    background-color: #0F766E;
    border-color: #0F766E;
}
</style>
""", unsafe_allow_html=True)


# ── Utility helpers ───────────────────────────────────────────────────────────

def _fmt_dollar(value: float, compact: bool = False) -> str:
    """Format *value* as a currency string.

    compact=True  → $1.2M / $450K  (for tight spaces like KPI cards)
    compact=False → $1,234,567      (for tables and tooltips)
    """
    if pd.isna(value):
        return "—"
    sign = "-" if value < 0 else ""
    abs_v = abs(value)
    if compact:
        if abs_v >= 1_000_000:
            return f"{sign}${abs_v / 1_000_000:.1f}M"
        if abs_v >= 1_000:
            return f"{sign}${abs_v / 1_000:.0f}K"
    return f"${value:,.0f}"


def _fmt_pct(value: float) -> str:
    """Format *value* as a signed percentage string, e.g. '+12.3%'."""
    if pd.isna(value):
        return "—"
    return f"{value:+.1f}%"


def _recompute_material(df: pd.DataFrame, mat_pct: float, mat_abs: float) -> pd.DataFrame:
    """
    Replace the Material column using user-specified thresholds from the sidebar sliders.

    variance_engine.compute_variances() uses the config-file defaults for materiality.
    This function overrides the Material column in-place on a copy so that slider
    changes are reflected without re-running the entire compute pipeline.

    Parameters
    ----------
    df      : enriched DataFrame (output of compute_variances)
    mat_pct : decimal threshold, e.g. 0.05 for 5%
    mat_abs : absolute dollar threshold, e.g. 50000
    """
    out = df.copy()
    abs_pct    = out["Variance (%)"].abs()
    abs_dollar = out["Variance ($)"].abs()
    out["Material"] = (abs_pct >= mat_pct * 100) | (abs_dollar >= mat_abs)
    return out


# ── KPI card HTML ─────────────────────────────────────────────────────────────

def _kpi_card(label: str, value: str, css_class: str = "", subtitle: str = "") -> str:
    """Return the HTML string for a single styled KPI card."""
    value_cls = f"kpi-value {css_class}".strip()
    sub = f'<div class="kpi-subtitle">{subtitle}</div>' if subtitle else ""
    return (
        f'<div class="kpi-card">'
        f'  <div class="kpi-label">{label}</div>'
        f'  <div class="{value_cls}">{value}</div>'
        f'  {sub}'
        f'</div>'
    )


# generate_commentary and answer_question are imported from commentary.py


# ── PowerPoint export ─────────────────────────────────────────────────────────

def _rgb(hex_color: str) -> RGBColor:
    """Convert a CSS hex string to a python-pptx RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_text_box(
    slide,
    text: str,
    left: float, top: float, width: float, height: float,
    font_size: int = 12,
    bold: bool = False,
    color: RGBColor | None = None,
    align=PP_ALIGN.LEFT,
) -> None:
    """Helper: add a formatted text box to *slide* using inch coordinates."""
    color = color or _rgb("#1E293B")
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.text = text
    para.alignment = align
    run = para.runs[0]
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.color.rgb = color
    run.font.name      = "Calibri"


def _chart_image(fig) -> BytesIO | None:
    """
    Convert a Plotly figure to a PNG BytesIO buffer via kaleido.

    Returns None (graceful degradation) if kaleido is not installed,
    so chart slides fall back to a text placeholder rather than crashing.
    """
    try:
        raw = fig.to_image(format="png", width=1100, height=520, scale=2)
        buf = BytesIO(raw)
        buf.seek(0)
        return buf
    except Exception:
        return None


def build_pptx(
    enriched_df: pd.DataFrame,
    cat_df: pd.DataFrame,
    waterfall_df: pd.DataFrame,
    period: str,
    commentary: str | None = None,
) -> BytesIO:
    """
    Assemble a PowerPoint report and return it as a BytesIO buffer.

    Slide structure
    ---------------
    1. Title slide  — period, headline metrics, byline
    2. Waterfall    — budget-to-actual bridge chart image
    3. By Category  — grouped bar chart image
    4. Material variances table — formatted with alternating row shading
    5. AI Commentary (only if *commentary* is provided)

    Chart images require the ``kaleido`` package (``pip install kaleido``).
    When kaleido is absent the image slides render a "install kaleido" notice
    instead of crashing.  The table and commentary slides are unaffected.

    Parameters
    ----------
    enriched_df  : output of compute_variances()
    cat_df       : output of summarize_by_category()
    waterfall_df : output of build_waterfall_data()
    period       : reporting period string used in slide titles
    commentary   : Claude commentary text, or None to omit the last slide
    """
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]   # fully blank layout

    SLATE  = _rgb("#1E293B")
    TEAL   = _rgb("#0D9488")
    LIGHT  = _rgb("#F8FAFC")
    MUTED  = _rgb("#94A3B8")
    WHITE  = _rgb("#FFFFFF")

    def _kaleido_notice(slide):
        _add_text_box(
            slide,
            "Chart image requires kaleido  →  pip install kaleido",
            0.4, 3.0, 12.5, 0.7,
            font_size=11, color=MUTED, align=PP_ALIGN.CENTER,
        )

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = _rgb("#0F172A")

    # Teal accent bar at bottom
    bar = sl.shapes.add_shape(1, Inches(0), Inches(6.85), Inches(13.33), Inches(0.65))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    _add_text_box(sl, "📊  FP&A Variance Copilot",
                  0.6, 1.4, 12.0, 1.0, font_size=30, bold=True, color=_rgb("#F1F5F9"))
    _add_text_box(sl, f"Variance Analysis Report  ·  {period}",
                  0.6, 2.6, 10.0, 0.7, font_size=17, color=MUTED)

    total_budget = enriched_df["Budget"].sum()
    total_actual = enriched_df["Actual"].sum()
    net_var      = total_actual - total_budget
    _add_text_box(
        sl,
        f"Budget {_fmt_dollar(total_budget, compact=True)}   ·   "
        f"Actual {_fmt_dollar(total_actual, compact=True)}   ·   "
        f"Net Variance {_fmt_dollar(net_var, compact=True)}",
        0.6, 3.5, 12.0, 0.6, font_size=13, color=_rgb("#CBD5E1"),
    )
    _add_text_box(sl, "Built by Priyanka Chandramohan",
                  0.6, 7.0, 6.0, 0.35, font_size=9, color=_rgb("#475569"))

    # ── Slide 2: Waterfall ────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _add_text_box(sl, f"Budget-to-Actual Bridge  ·  {period}",
                  0.4, 0.2, 12.5, 0.55, font_size=14, bold=True, color=SLATE)
    img = _chart_image(plot_waterfall(waterfall_df, title=""))
    if img:
        sl.shapes.add_picture(img, Inches(0.4), Inches(0.9), Inches(12.5), Inches(5.9))
    else:
        _kaleido_notice(sl)

    # ── Slide 3: Budget vs Actual by Category ─────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _add_text_box(sl, f"Budget vs. Actual by Category  ·  {period}",
                  0.4, 0.2, 12.5, 0.55, font_size=14, bold=True, color=SLATE)
    img = _chart_image(plot_budget_vs_actual(cat_df, title=""))
    if img:
        sl.shapes.add_picture(img, Inches(0.4), Inches(0.9), Inches(12.5), Inches(5.9))
    else:
        _kaleido_notice(sl)

    # ── Slide 4: Material variances table ─────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _add_text_box(sl, f"Material Variances  ·  {period}",
                  0.4, 0.2, 12.5, 0.55, font_size=14, bold=True, color=SLATE)

    material = get_material_items(enriched_df)
    tbl_cols  = ["Line Item", "Category", "Budget", "Actual", "Variance ($)", "Variance (%)", "Severity"]
    tbl_data  = material[[c for c in tbl_cols if c in material.columns]]
    n_rows    = len(tbl_data) + 1   # +1 for header

    tbl = sl.shapes.add_table(
        n_rows, len(tbl_cols),
        Inches(0.4), Inches(0.9),
        Inches(12.5), Inches(min(5.8, 0.38 * n_rows + 0.3)),
    ).table

    # Header
    for ci, col_name in enumerate(tbl_cols):
        cell = tbl.cell(0, ci)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = SLATE
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(8); run.font.bold = True
        run.font.color.rgb = WHITE; run.font.name = "Calibri"

    # Data rows
    for ri, (_, row) in enumerate(tbl_data.iterrows(), start=1):
        for ci, col_name in enumerate(tbl_cols):
            cell = tbl.cell(ri, ci)
            val  = row[col_name]
            if col_name in ("Budget", "Actual", "Variance ($)"):
                cell.text = _fmt_dollar(val)
            elif col_name == "Variance (%)":
                cell.text = _fmt_pct(val)
            else:
                cell.text = str(val)
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(8); run.font.name = "Calibri"
            run.font.color.rgb = SLATE
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb("#F8FAFC")

    # ── Slide 5: AI Commentary ────────────────────────────────────────────────
    if commentary:
        sl = prs.slides.add_slide(blank)
        _add_text_box(sl, f"AI Commentary  ·  {period}",
                      0.4, 0.2, 12.5, 0.55, font_size=14, bold=True, color=SLATE)
        # Left teal accent bar
        accent = sl.shapes.add_shape(
            1, Inches(0.4), Inches(0.9), Inches(0.06), Inches(5.8)
        )
        accent.fill.solid(); accent.fill.fore_color.rgb = TEAL
        accent.line.fill.background()

        txb = sl.shapes.add_textbox(Inches(0.6), Inches(0.9), Inches(12.3), Inches(5.8))
        tf  = txb.text_frame; tf.word_wrap = True
        first = True
        for para_text in commentary.split("\n"):
            para_text = para_text.strip()
            if not para_text:
                continue
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = para_text
            run = p.runs[0]
            run.font.size = Pt(10); run.font.name = "Calibri"
            run.font.color.rgb = _rgb("#334155")

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ── State 1: Landing page ─────────────────────────────────────────────────────

def render_landing() -> None:
    """Render the welcome screen shown when no file has been uploaded."""
    st.markdown("<br>", unsafe_allow_html=True)
    st.markdown(
        "<h1 style='text-align:center;color:#1E293B;"
        "font-family:Calibri,sans-serif;font-size:2rem;font-weight:700;"
        "line-height:1.25;'>Your CFO-grade variance analyst,<br>always on call.</h1>",
        unsafe_allow_html=True,
    )
    st.markdown(
        "<p style='text-align:center;color:#64748B;font-size:0.975rem;"
        "font-family:Calibri,sans-serif;max-width:600px;margin:12px auto 40px;'>"
        "Drop in your budget vs. actuals CSV and get an instant waterfall bridge, "
        "severity breakdown, and CFO-ready commentary — in seconds.</p>",
        unsafe_allow_html=True,
    )

    g1, g2, g3, g4, g5 = st.columns([6, 1, 6, 1, 6], gap="small")
    steps = [
        (g1, "📁", "Upload",
         "Drop in any P&L CSV with Line Item, Category, Budget, and Actual columns. "
         "Use the sample file to explore before connecting your own data."),
        (g3, "📊", "Analyze",
         "Instant waterfall bridge from budget to actual, severity heat map, "
         "and category roll-up — material items surface automatically."),
        (g5, "💬", "Ask & Export",
         "Add your Anthropic API key for CFO-level narrative from Claude. "
         "Download a polished PowerPoint deck in one click."),
    ]
    connectors = [g2, g4]

    for col, icon, heading, body in steps:
        with col:
            st.markdown(
                f'<div class="landing-step">'
                f'<div class="step-icon">{icon}</div>'
                f'<h3>{heading}</h3>'
                f'<p>{body}</p>'
                f'</div>',
                unsafe_allow_html=True,
            )
    for conn_col in connectors:
        with conn_col:
            st.markdown(
                '<div class="landing-connector">→</div>',
                unsafe_allow_html=True,
            )

    st.markdown("<br><br>", unsafe_allow_html=True)
    st.markdown(
        "<p style='text-align:center;color:#94A3B8;font-size:0.78rem;"
        "font-family:Calibri,sans-serif;'>← Upload your CSV in the sidebar to get started</p>",
        unsafe_allow_html=True,
    )


# ── State 2 helpers ───────────────────────────────────────────────────────────

def render_kpi_cards(enriched_df: pd.DataFrame, period: str) -> None:
    """Render four KPI metric cards across the top of the analysis area."""
    total_budget   = enriched_df["Budget"].sum()
    total_actual   = enriched_df["Actual"].sum()
    net_var        = total_actual - total_budget
    material_mask  = enriched_df["Material"]
    material_count = int(material_mask.sum())
    unfav_count    = int((material_mask & ~enriched_df["Favorable"]).sum())

    net_css  = "favorable" if net_var >= 0 else "unfavorable"
    net_sign = "+" if net_var >= 0 else ""

    c1, c2, c3, c4 = st.columns(4, gap="medium")
    cards = [
        (c1, "Total Budget",  _fmt_dollar(total_budget, compact=True), "",       period),
        (c2, "Total Actual",  _fmt_dollar(total_actual, compact=True), "",       period),
        (c3, "Net Variance",  f"{net_sign}{_fmt_dollar(net_var, compact=True)}", net_css, "vs. Budget"),
        (c4, "Material Items", str(material_count),
             "unfavorable" if unfav_count > material_count / 2 else "",
             f"{unfav_count} unfavorable"),
    ]
    for col, label, value, css_class, subtitle in cards:
        with col:
            st.markdown(_kpi_card(label, value, css_class, subtitle), unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)


def render_data_table(enriched_df: pd.DataFrame) -> None:
    """
    Render all line items as a formatted, styled dataframe.

    Variance ($) and the Favorable column are colour-coded green/red.
    Dollar columns use Streamlit's NumberColumn formatter so they render
    with $ and commas without requiring a pandas Styler round-trip.
    """
    display_cols = [
        "Line Item", "Category", "Budget", "Actual",
        "Variance ($)", "Variance (%)", "Favorable", "Severity", "Material",
    ]
    df = enriched_df[[c for c in display_cols if c in enriched_df.columns]].copy()

    def _colour_row(row):
        """Apply green/red to Variance ($) and Favorable cells."""
        styles  = [""] * len(row)
        is_fav  = bool(row.get("Favorable", True))
        colour  = "#16A34A" if is_fav else "#DC2626"
        for col_name, style in [
            ("Variance ($)", f"color:{colour};"),
            ("Favorable",    f"color:{colour};font-weight:600;"),
        ]:
            if col_name in row.index:
                styles[row.index.get_loc(col_name)] = style
        return styles

    styled = df.style.apply(_colour_row, axis=1)

    st.dataframe(
        styled,
        use_container_width=True,
        column_config={
            "Budget":       st.column_config.NumberColumn(format="$%d"),
            "Actual":       st.column_config.NumberColumn(format="$%d"),
            "Variance ($)": st.column_config.NumberColumn(format="$%d"),
            "Variance (%)": st.column_config.NumberColumn(format="%.1f%%"),
            "Favorable":    st.column_config.CheckboxColumn(disabled=True),
            "Material":     st.column_config.CheckboxColumn(disabled=True),
        },
        hide_index=True,
        height=440,
    )


def render_analysis(
    enriched_df: pd.DataFrame,
    cat_df: pd.DataFrame,
    waterfall_df: pd.DataFrame,
    period: str,
) -> None:
    """Render KPI cards and the four-tab chart / table area."""
    render_kpi_cards(enriched_df, period)

    tab_wf, tab_cat, tab_sev, tab_tbl = st.tabs([
        "📉  Waterfall",
        "📊  By Category",
        "🔴  Severity",
        "📋  Data Table",
    ])

    with tab_wf:
        st.plotly_chart(
            plot_waterfall(waterfall_df, title=f"{period} — Budget-to-Actual Bridge"),
            use_container_width=True,
        )

    with tab_cat:
        st.plotly_chart(
            plot_budget_vs_actual(cat_df, title=f"{period} — Budget vs. Actual by Category"),
            use_container_width=True,
        )

    with tab_sev:
        col_a, col_b = st.columns(2, gap="large")
        with col_a:
            st.plotly_chart(
                plot_severity_donut(enriched_df, title="All Variances by Severity"),
                use_container_width=True,
            )
        with col_b:
            st.plotly_chart(
                plot_severity_donut(enriched_df, title="Material Variances Only",
                                    material_only=True),
                use_container_width=True,
            )

    with tab_tbl:
        render_data_table(enriched_df)


# ── State 3: Commentary section ───────────────────────────────────────────────

def render_commentary(text: str) -> None:
    """Render the Claude commentary in a styled box below the charts."""
    st.markdown("---")
    st.markdown(
        "<h3 style='font-family:Calibri,sans-serif;color:#1E293B;"
        "font-size:1.1rem;font-weight:700;margin-bottom:12px;'>💬 AI Commentary</h3>",
        unsafe_allow_html=True,
    )
    # Escape HTML entities but preserve newlines as <br> for the box
    safe_text = (
        text
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace("\n", "<br>")
    )
    st.markdown(f'<div class="commentary-box">{safe_text}</div>', unsafe_allow_html=True)


# ── Sidebar ───────────────────────────────────────────────────────────────────

def render_sidebar() -> dict:
    """
    Render the sidebar controls and return a dict of current values.

    Returns
    -------
    dict with keys:
        uploaded_file : UploadedFile | None
        period        : str
        mat_pct       : float  (decimal, e.g. 0.05)
        mat_abs       : float  (dollars, e.g. 50000)
        tone          : str    (one of TONES)
        api_key       : str    (empty string if not provided)
    """
    with st.sidebar:
        st.markdown("## 📊 FP&A Variance Copilot")
        st.markdown("---")

        uploaded_file = st.file_uploader(
            "Upload CSV",
            type=["csv"],
            help="Required columns: Line Item, Category, Budget, Actual",
        )

        period = st.text_input("Reporting Period", value="Q4 2024")

        st.markdown("---")
        st.markdown(
            "<p style='font-size:0.70rem;color:#94A3B8;letter-spacing:0.06em;"
            "text-transform:uppercase;margin-bottom:4px;'>Commentary Tone</p>",
            unsafe_allow_html=True,
        )
        tone = st.radio(
            "Commentary Tone",
            options=TONES,
            index=TONES.index(DEFAULT_TONE),
            label_visibility="collapsed",
            help=(
                "Board Presentation: polished, strategic grouping, recommendations.\n"
                "Internal Team Update: direct, action items, under 300 words.\n"
                "Investor Call Prep: formal, beat/miss framing, Q&A prep."
            ),
        )

        st.markdown("---")
        st.markdown(
            "<p style='font-size:0.70rem;color:#94A3B8;letter-spacing:0.06em;"
            "text-transform:uppercase;margin-bottom:4px;'>Materiality Thresholds</p>",
            unsafe_allow_html=True,
        )
        mat_pct = st.slider(
            "% Threshold", min_value=1, max_value=25, value=5,
            help="A variance is material if |Var%| ≥ this value",
        )
        mat_abs = st.slider(
            "$ Threshold (K)", min_value=10, max_value=500, value=50, step=10,
            help="A variance is material if |Var$| ≥ this × $1,000",
        )

        st.markdown("---")
        api_key = st.text_input(
            "Anthropic API Key",
            type="password",
            placeholder="sk-ant-…  (optional)",
            help="Required for AI commentary and querying. Key is not stored anywhere.",
        )

        # Push byline to the bottom of the sidebar
        st.markdown("<br>" * 6, unsafe_allow_html=True)
        st.markdown(
            "<p style='font-size:0.70rem;color:#475569;text-align:center;"
            "font-family:Calibri,sans-serif;line-height:1.6;'>"
            "Built by<br>"
            "<strong style='color:#94A3B8;font-size:0.78rem;'>"
            "Priyanka Chandramohan</strong></p>",
            unsafe_allow_html=True,
        )

    return {
        "uploaded_file": uploaded_file,
        "period":        period,
        "mat_pct":       mat_pct / 100,
        "mat_abs":       mat_abs * 1_000,
        "tone":          tone,
        "api_key":       api_key.strip() if api_key else "",
    }


# ── Chat interface (NL querying) ──────────────────────────────────────────────

def render_chat_interface(
    enriched_df: pd.DataFrame,
    cat_df: pd.DataFrame,
    period: str,
    api_key: str,
) -> None:
    """
    Render a chat-style interface for natural-language variance queries.

    When an API key is present, each user message is sent to Claude with the
    full variance context and the accumulated conversation history.  Responses
    are streamed into st.chat_message so multi-turn follow-ups work naturally.

    When no API key is available, a prompt is shown directing the user to the
    sidebar — the chat input is still rendered so the UX is consistent.

    Session-state key: "chat_history" — list of {"role", "content"} dicts.
    The history is preserved until a new file is uploaded or thresholds change.
    """
    st.markdown("---")
    st.markdown(
        "<h3 style='font-family:Calibri,sans-serif;color:#1E293B;"
        "font-size:1.1rem;font-weight:700;margin-bottom:4px;'>💬 Ask About Your Variances</h3>",
        unsafe_allow_html=True,
    )

    if not api_key:
        st.info(
            "Add your Anthropic API key in the sidebar to ask questions about your variances.",
            icon="🔑",
        )
    else:
        st.markdown(
            "<p style='font-size:0.8rem;color:#94A3B8;font-family:Calibri,sans-serif;"
            "margin-bottom:12px;'>"
            "Try: \"Why did we miss revenue?\" · \"What's driving the cloud overrun?\" · "
            "\"Summarise this for my CFO in 2 sentences\"</p>",
            unsafe_allow_html=True,
        )

    history: list[dict] = st.session_state.get("chat_history", [])

    # Render existing conversation
    for msg in history:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    # Chat input (always rendered; disabled path handled below)
    if prompt := st.chat_input("Ask a question about your variances…"):
        # Immediately render the user bubble and save to history
        with st.chat_message("user"):
            st.markdown(prompt)
        history.append({"role": "user", "content": prompt})

        # Generate response
        with st.chat_message("assistant"):
            if not api_key:
                reply = (
                    "Add your Anthropic API key in the sidebar to ask questions "
                    "about your variances."
                )
                st.markdown(reply)
            else:
                with st.spinner(""):
                    try:
                        reply = answer_question(
                            enriched_df=enriched_df,
                            cat_df=cat_df,
                            period=period,
                            question=prompt,
                            chat_history=history[:-1],   # exclude the message we just appended
                            api_key=api_key,
                            commentary=st.session_state.get("commentary"),
                        )
                    except anthropic.AuthenticationError:
                        reply = "Invalid API key — please check the key in the sidebar."
                    except anthropic.APIConnectionError:
                        reply = "Could not reach the Anthropic API. Check your internet connection."
                    except Exception as exc:
                        reply = f"Error: {exc}"
                st.markdown(reply)

        history.append({"role": "assistant", "content": reply})
        st.session_state["chat_history"] = history
        st.rerun()


# ── Main ──────────────────────────────────────────────────────────────────────

def main() -> None:
    inputs = render_sidebar()

    uploaded_file = inputs["uploaded_file"]
    period        = inputs["period"]
    mat_pct       = inputs["mat_pct"]
    mat_abs       = inputs["mat_abs"]
    tone          = inputs["tone"]
    api_key       = inputs["api_key"]

    # ── State 1: no file ─────────────────────────────────────────────────────
    if uploaded_file is None:
        render_landing()
        return

    # ── Load & validate ───────────────────────────────────────────────────────
    try:
        df_raw = pd.read_csv(uploaded_file)
    except Exception as exc:
        st.error(f"Could not read the uploaded file: {exc}")
        return

    ok, msg = validate_csv(df_raw)
    if not ok:
        st.error(f"**Validation error:** {msg}")
        return

    # Cache-bust key: file identity + materiality settings
    # If any of these change, recompute from scratch and clear commentary.
    file_key = f"{uploaded_file.name}|{uploaded_file.size}|{mat_pct}|{mat_abs}"

    if st.session_state.get("_file_key") != file_key:
        # New file or changed thresholds — run the full pipeline
        enriched_df  = compute_variances(df_raw)
        enriched_df  = _recompute_material(enriched_df, mat_pct, mat_abs)
        cat_df       = summarize_by_category(enriched_df)
        waterfall_df = build_waterfall_data(enriched_df)

        st.session_state["_file_key"]    = file_key
        st.session_state["enriched_df"]  = enriched_df
        st.session_state["cat_df"]       = cat_df
        st.session_state["waterfall_df"] = waterfall_df
        st.session_state["commentary"]   = None   # reset on new file / threshold change
        st.session_state["chat_history"] = []     # reset conversation on new file

    enriched_df  = st.session_state["enriched_df"]
    cat_df       = st.session_state["cat_df"]
    waterfall_df = st.session_state["waterfall_df"]

    # ── State 2: analysis ─────────────────────────────────────────────────────
    render_analysis(enriched_df, cat_df, waterfall_df, period)

    # ── Action row ────────────────────────────────────────────────────────────
    st.markdown("---")
    btn_col, dl_col, _ = st.columns([2, 2, 5], gap="small")

    with btn_col:
        # Show which tone will be used so users understand the button
        btn_label = f"✨ Generate  ·  {tone.split()[0]}"
        generate_clicked = st.button(
            btn_label,
            type="primary",
            use_container_width=True,
            help=(
                f"Generate {tone} commentary. "
                "No API key? Click anyway for a rule-based version."
            ),
        )

    with dl_col:
        pptx_buf = build_pptx(
            enriched_df, cat_df, waterfall_df, period,
            commentary=st.session_state.get("commentary"),
        )
        st.download_button(
            label="📥 Download PowerPoint",
            data=pptx_buf,
            file_name=f"variance_report_{period.replace(' ', '_')}.pptx",
            mime=(
                "application/vnd.openxmlformats-officedocument"
                ".presentationml.presentation"
            ),
            use_container_width=True,
        )

    # ── Commentary generation ─────────────────────────────────────────────────
    if generate_clicked:
        spinner_msg = (
            f"Claude is preparing {tone} commentary…"
            if api_key else
            f"Generating rule-based {tone} commentary…"
        )
        with st.spinner(spinner_msg):
            try:
                text = generate_commentary(
                    enriched_df, cat_df, period,
                    tone=tone,
                    api_key=api_key,
                )
                st.session_state["commentary"] = text
                st.session_state["chat_history"] = []   # fresh chat for new commentary
                st.rerun()
            except anthropic.AuthenticationError:
                st.error("Invalid API key. Double-check your Anthropic key and try again.")
            except anthropic.APIConnectionError:
                st.error("Could not reach the Anthropic API. Check your internet connection.")
            except Exception as exc:
                st.error(f"Commentary generation failed: {exc}")

    # ── State 3: commentary + chat ────────────────────────────────────────────
    if st.session_state.get("commentary"):
        render_commentary(st.session_state["commentary"])
        render_chat_interface(enriched_df, cat_df, period, api_key)


if __name__ == "__main__":
    main()
